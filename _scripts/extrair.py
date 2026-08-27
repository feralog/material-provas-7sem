#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
extrair.py — extrai texto E imagens de uma aula (PDF / PPTX / DOCX / TXT / MD).

Gera dois arquivos por aula:
    conteudo_[tema].md     texto, com marcadores <!-- IMG:pgNN --> no fim de cada pagina
    imagens_[tema].json    { "pgNN": "<base64 JPEG>" } — uma entrada por pagina

O .md e a fonte de verdade para escrever quiz e resumo.
O .json e o acervo de imagens: o montar_quiz.py embute no HTML APENAS as que forem
referenciadas, entao extrair tudo aqui nao pesa no arquivo final.

Uso:
    python _scripts/extrair.py "PED/Aula 4 - Tireoide.pdf"
    python _scripts/extrair.py "PED" --por-arquivo      # 1 par de arquivos por aula
    python _scripts/extrair.py "PED" --sem-imagens      # so o texto
    python _scripts/extrair.py "PED" --largura 1600     # imagens maiores
"""

import argparse
import base64
import datetime
import hashlib
import json
import logging
import re
import sys
import unicodedata
from pathlib import Path

# pdfminer reclama de FontBBox em quase todo slide exportado — ruido puro
logging.getLogger("pdfminer").setLevel(logging.ERROR)
logging.getLogger("pdfplumber").setLevel(logging.ERROR)

EXTS = {".pdf", ".pptx", ".docx", ".txt", ".md"}
LARGURA = 1200      # px da imagem renderizada
QUALIDADE = 82      # qualidade JPEG


def snake(s: str) -> str:
    s = unicodedata.normalize("NFKD", s).encode("ascii", "ignore").decode()
    s = re.sub(r"[^\w\s-]", "", s).strip().lower()
    s = re.sub(r"[\s\-]+", "_", s)
    return re.sub(r"_+", "_", s).strip("_") or "conteudo"


def titulo_limpo(nome: str) -> str:
    """Remove sufixos de exportacao do nome do arquivo: '_260819_133948', '.pptx'."""
    t = re.sub(r"\.(pptx|docx|pdf)$", "", nome, flags=re.I)
    t = re.sub(r"[_\- ]*\d{6}[_\- ]\d{6}$", "", t)
    return re.sub(r"\s{2,}", " ", t).strip(" -_")


def normalizar_imagem(blob: bytes, largura: int):
    """Reduz a imagem para `largura` px e reencoda em JPEG. Devolve None se não der.

    Duas razões para existir:
    1. Imagens embutidas em PPTX vêm em resolução original (4000 px, 1,3 MB) para
       uma coluna de 720 px — sem reduzir, o HTML final passa de 10 MB.
    2. O HTML embute tudo como `data:image/jpeg`. Formato vetorial (EMF/WMF), que o
       PowerPoint usa para diagramas colados, apareceria como figura quebrada.
       Melhor descartar na extração do que guardar um blob que não dá para exibir.
    """
    try:
        import fitz
        doc = fitz.open(stream=blob, filetype="image")
        page = doc[0]
        zoom = min(1.0, largura / page.rect.width) if page.rect.width else 1.0
        pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
        if pix.alpha:                       # achata transparência sobre branco
            pix = fitz.Pixmap(fitz.csRGB, pix)
        out = pix.tobytes("jpg", jpg_quality=QUALIDADE)
        doc.close()
    except Exception:
        return None
    if not out or out[:3] != b"\xff\xd8\xff":
        return None
    # se o original ja for um JPEG menor, fica com ele
    return blob if (blob[:3] == b"\xff\xd8\xff" and len(blob) < len(out)) else out


def limpar(texto: str) -> str:
    """Normaliza espacos e remove linhas de ruido tipicas de slide."""
    linhas = []
    for ln in texto.splitlines():
        ln = ln.replace("\xa0", " ").rstrip()
        ln = re.sub(r"[ \t]{2,}", " ", ln)
        if not ln.strip():
            linhas.append("")
            continue
        if re.fullmatch(r"\d{1,3}", ln.strip()):
            continue
        linhas.append(ln)
    txt = "\n".join(linhas)
    return re.sub(r"\n{3,}", "\n\n", txt).strip()


# ── extratores ───────────────────────────────────────────────────────────────

def extrair_pdf(path: Path, com_imagens: bool, largura: int):
    """Retorna (markdown, imagens{key:b64}, avisos)."""
    import fitz  # PyMuPDF

    texto_por_pagina, vazias = {}, []
    try:
        import pdfplumber
        with pdfplumber.open(str(path)) as pdf:
            total = len(pdf.pages)
            for i, page in enumerate(pdf.pages, 1):
                texto_por_pagina[i] = limpar(page.extract_text() or "")
    except Exception as e:
        texto_por_pagina = {}
        doc = fitz.open(str(path))
        total = doc.page_count
        for i, page in enumerate(doc, 1):
            texto_por_pagina[i] = limpar(page.get_text())
        doc.close()
        _ = e

    imagens = {}
    doc = fitz.open(str(path))
    total = doc.page_count
    partes = []
    for i in range(1, total + 1):
        key = f"pg{i:02d}"
        t = texto_por_pagina.get(i, "")
        if not t:
            vazias.append(i)
        if com_imagens:
            page = doc[i - 1]
            zoom = largura / page.rect.width
            pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
            imagens[key] = base64.b64encode(pix.tobytes("jpg", jpg_quality=QUALIDADE)).decode()
        cabec = f"### Página {i}"
        corpo = t if t else "_(slide sem texto — o conteúdo está só na imagem)_"
        marca = f"\n\n<!-- IMG:{key} -->" if com_imagens else ""
        partes.append(f"{cabec}\n\n{corpo}{marca}")
    doc.close()

    avisos = []
    if vazias:
        avisos.append(f"{path.name}: {len(vazias)}/{total} páginas sem texto nativo "
                      f"(o conteúdo está na imagem) -> {vazias[:20]}")
    return "\n\n".join(partes), imagens, avisos


def extrair_pptx(path: Path, com_imagens: bool, largura: int):
    try:
        from pptx import Presentation
    except ImportError:
        return "", {}, [f"{path.name}: python-pptx não instalado (pip install python-pptx)"]

    EMU_POL = 914400          # 1 polegada
    MIN_LADO = EMU_POL // 2   # imagem menor que meia polegada = decoracao
    prs = Presentation(str(path))

    def texto_frame(tf):
        """Preserva os niveis de bullet como lista markdown."""
        linhas = []
        for p in tf.paragraphs:
            t = "".join(r.text for r in p.runs).strip() or p.text.strip()
            if not t:
                continue
            linhas.append(("  " * min(p.level, 3) + "- " + t) if p.level or len(tf.paragraphs) > 1 else t)
        return "\n".join(linhas)

    def tabela_md(tb):
        linhas = [[c.text.strip().replace("\n", " ") for c in row.cells] for row in tb.rows]
        if not linhas:
            return ""
        larg = len(linhas[0])
        out = ["| " + " | ".join(linhas[0]) + " |",
               "| " + " | ".join(["---"] * larg) + " |"]
        for l in linhas[1:]:
            out.append("| " + " | ".join(l) + " |")
        return "\n".join(out)

    def anda(shapes, saida, figuras):
        """Percorre as shapes (inclusive dentro de grupos) em ordem de leitura."""
        try:
            ordenadas = sorted(shapes, key=lambda s: (s.top or 0, s.left or 0))
        except Exception:
            ordenadas = list(shapes)
        for sh in ordenadas:
            if sh.shape_type == 6:  # GROUP
                anda(sh.shapes, saida, figuras)
                continue
            if getattr(sh, "has_table", False):
                md = tabela_md(sh.table)
                if md:
                    saida.append(md)
                continue
            if sh.shape_type == 13:  # PICTURE
                if (sh.width or 0) >= MIN_LADO and (sh.height or 0) >= MIN_LADO:
                    try:
                        figuras.append(sh.image.blob)
                    except Exception:
                        pass
                continue
            if sh.has_text_frame:
                t = texto_frame(sh.text_frame)
                if t:
                    saida.append(t)

    # 1a passada: junta tudo, contando repeticao de imagem (logo aparece em todo slide)
    slides, contagem = [], {}
    for i, slide in enumerate(prs.slides, 1):
        saida, figuras = [], []
        anda(slide.shapes, saida, figuras)
        notas = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            notas = slide.notes_slide.notes_text_frame.text.strip()
        digests = []
        for blob in figuras:
            h = hashlib.md5(blob).hexdigest()
            contagem[h] = contagem.get(h, 0) + 1
            digests.append((h, blob))
        slides.append((i, saida, notas, digests))

    repetidas = {h for h, n in contagem.items() if n >= 4}   # marca d'agua / logo

    partes, imagens, vistos, descartadas = [], {}, {}, 0
    for i, saida, notas, digests in slides:
        marcas = []
        if com_imagens:
            for h, blob in digests:
                if h in repetidas:
                    continue
                if h in vistos:
                    marcas.append(vistos[h])
                    continue
                jpg = normalizar_imagem(blob, largura)
                if jpg is None:          # vetorial (EMF/WMF) — não exibível no HTML
                    descartadas += 1
                    continue
                key = f"s{i:02d}i{len([k for k in imagens if k.startswith(f's{i:02d}i')]) + 1}"
                imagens[key] = base64.b64encode(jpg).decode()
                vistos[h] = key
                marcas.append(key)
        corpo = limpar("\n\n".join(saida))
        if notas:
            corpo += ("\n\n" if corpo else "") + "**Notas do apresentador:** " + limpar(notas)
        if not corpo and not marcas:
            continue
        bloco = f"### Slide {i}\n\n{corpo or '_(slide sem texto — o conteúdo está na imagem)_'}"
        bloco += "".join(f"\n\n<!-- IMG:{k} -->" for k in marcas)
        partes.append(bloco)

    avisos = []
    if repetidas:
        avisos.append(f"{path.name}: {len(repetidas)} imagem(ns) repetida(s) em 4+ slides "
                      f"descartada(s) como logo/marca d'água")
    if descartadas:
        avisos.append(f"{path.name}: {descartadas} imagem(ns) vetorial(is) (EMF/WMF) "
                      f"descartada(s) — não são exibíveis no HTML")
    return "\n\n".join(partes), imagens, avisos


def extrair_docx(path: Path, com_imagens: bool, largura: int):
    try:
        import docx
    except ImportError:
        return "", {}, [f"{path.name}: python-docx não instalado (pip install python-docx)"]
    d = docx.Document(str(path))
    buf = [p.text for p in d.paragraphs]
    for tb in d.tables:
        for row in tb.rows:
            buf.append(" | ".join(c.text.strip() for c in row.cells))
    return limpar("\n".join(buf)), {}, []


def extrair_txt(path: Path, com_imagens: bool, largura: int):
    return limpar(path.read_text(encoding="utf-8", errors="replace")), {}, []


EXTRATORES = {".pdf": extrair_pdf, ".pptx": extrair_pptx, ".docx": extrair_docx,
              ".txt": extrair_txt, ".md": extrair_txt}


# ── montagem ─────────────────────────────────────────────────────────────────

def processar(arquivos, titulo, origem, com_imagens, largura):
    corpo, avisos, processados, imagens = [], [], [], {}
    multi = len(arquivos) > 1
    for f in arquivos:
        fn = EXTRATORES.get(f.suffix.lower())
        if not fn:
            continue
        print(f"  lendo {f.name} ...", flush=True)
        txt, imgs, av = fn(f, com_imagens, largura)
        avisos += av
        if not txt.strip():
            avisos.append(f"{f.name}: nenhum texto extraído")
            continue
        # PDF gera chaves pgNN e PPTX gera sNNiK, entao em geral nao colidem.
        # So prefixa quando houver colisao de verdade (ex.: dois PDFs no mesmo material).
        if multi and any(k in imagens for k in imgs):
            pref = f"f{len(processados) + 1}"
            imgs = {f"{pref}_{k}": v for k, v in imgs.items()}
            txt = re.sub(r"<!-- IMG:(\w+) -->", rf"<!-- IMG:{pref}_\1 -->", txt)
        imagens.update(imgs)
        processados.append(f.name)
        # com vários arquivos o nome do stem pode repetir (mesma aula em .pptx e .pdf)
        rotulo = titulo_limpo(f.stem) + (f"  ·  {f.suffix.lstrip('.').upper()}" if multi else "")
        corpo.append(f"## {rotulo}\n\n{txt}")

    head = (
        "---\n"
        f"source: {origem}\n"
        f"files_processed: {processados}\n"
        f"date_extracted: {datetime.date.today().isoformat()}\n"
        f"subject: {titulo}\n"
        f"imagens: {len(imagens)} (chaves pgNN — ver imagens_[tema].json)\n"
        "---\n\n"
        f"# {titulo}\n\n"
    )
    return head + "\n\n---\n\n".join(corpo) + "\n", imagens, avisos


def main():
    ap = argparse.ArgumentParser(description="Extrai texto e imagens de aula")
    ap.add_argument("alvo", nargs="+",
                    help="arquivo(s) ou pasta da aula. Vários arquivos = um material só "
                         "(ex.: os slides .pptx + o handout .pdf da mesma aula)")
    ap.add_argument("--titulo", help="título do material (padrão: nome do 1º arquivo)")
    ap.add_argument("--saida", help="caminho do .md de saída")
    ap.add_argument("--por-arquivo", action="store_true",
                    help="gera um .md por arquivo da pasta")
    ap.add_argument("--sem-imagens", action="store_true", help="extrair só o texto")
    ap.add_argument("--largura", type=int, default=LARGURA, help=f"px da imagem (padrão {LARGURA})")
    args = ap.parse_args()

    alvos = [Path(a).expanduser().resolve() for a in args.alvo]
    for a in alvos:
        if not a.exists():
            sys.exit(f"ERRO: não encontrei {a}")
    com_imagens = not args.sem_imagens

    if len(alvos) > 1:
        # varios arquivos -> UM material so
        if any(a.is_dir() for a in alvos):
            sys.exit("ERRO: com vários alvos, todos precisam ser arquivos (não pastas)")
        titulo = args.titulo or titulo_limpo(alvos[0].stem)
        grupos = [(titulo, alvos, alvos[0].parent)]
    elif alvos[0].is_file():
        grupos = [(args.titulo or titulo_limpo(alvos[0].stem), [alvos[0]], alvos[0].parent)]
    else:
        alvo = alvos[0]
        arquivos = sorted(p for p in alvo.iterdir()
                          if p.is_file() and p.suffix.lower() in EXTS
                          and not p.name.startswith(("conteudo_", "resumo_")))
        if not arquivos:
            sys.exit(f"ERRO: nenhum arquivo suportado em {alvo}")
        grupos = ([(titulo_limpo(f.stem), [f], alvo) for f in arquivos]
                  if args.por_arquivo else [(args.titulo or alvo.name, arquivos, alvo)])
    alvo = alvos[0]

    todos_avisos = []
    for titulo, arquivos, pasta in grupos:
        print(f"\n>> {titulo}  ({len(arquivos)} arquivo(s))")
        md, imagens, avisos = processar(arquivos, titulo, str(alvo), com_imagens, args.largura)
        todos_avisos += avisos

        out = (Path(args.saida).expanduser().resolve()
               if args.saida and len(grupos) == 1
               else pasta / f"conteudo_{snake(titulo)}.md")
        out.parent.mkdir(parents=True, exist_ok=True)
        out.write_text(md, encoding="utf-8")
        print(f"   texto   -> {out.name}  ({len(md):,} chars)")

        if imagens:
            ij = out.parent / f"imagens_{snake(titulo)}.json"
            ij.write_text(json.dumps(imagens, ensure_ascii=False), encoding="utf-8")
            mb = sum(len(v) for v in imagens.values()) / 1024 / 1024
            print(f"   imagens -> {ij.name}  ({len(imagens)} páginas, {mb:.1f} MB em base64)")

    if todos_avisos:
        print("\nAVISOS:")
        for a in todos_avisos:
            print("  ! " + a)
        print("\n  Páginas sem texto: use a imagem correspondente (<!-- IMG:pgNN -->)\n"
              "  ou leia o PDF com a tool Read, que renderiza a página.")


if __name__ == "__main__":
    main()
